from ensurepip import bootstrap
import re
import sys
sys.path.append('c:\\Users\\long.pham\\Documents\\MDDPython')

from email.policy import default
from inspect import stack
from ipaddress import v4_int_to_packed
import math
from re import T
from turtle import left, pos
from xml.dom.pulldom import ErrorHandler
import win32com.client as w32
from metadata import Metadata
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import openpyxl
from pptx import Presentation
from pptx.util import Inches


#MDM = w32.Dispatch('MDM.Document')
#MDM.Open(r'metadatas/PVN2021249_F61_W1_2022v1.mdd')

#questions = ["InstanceID","_Year","_Month","_Quarter","_ResProvinces","_City_Groups","_Target","_Users","_Class","_EPI_OOP_Inject","_Table_1","_Table_2","_Table_8","_Table_9","_Table_10_Before","_Table_10","_Table_13","_Table_14","_Table_15","_Table_16","_Table_17","_Table_18_Before","_Table_18","_Table_19","_Table_20","_Table_21","_Table_3","_UA1","_UA2","_UA3a","_UA3b","_UB1","_UB2","_UB4","_UB5","_Q3","_S16","_S16b","_S16c","_Q17","_UD3","_OP5c","_OP5e","_OP5d"]

project_name = 'prj_tablo'
mdd_filename = 'metadatas/S22015945_DATA.mdd'
ddf_filename = 'metadatas/S22015945_DATA.ddf'

questions = ["Respondent.ID","_Q1","_Phase"]

#try:
m = Metadata(r'{}/{}'.format(project_name, mdd_filename), r'{}/{}'.format(project_name, ddf_filename), questions)

df = m.convertToDataFrame() 
df = df.set_index("Respondent.ID")

obj_skus_define = {
    '_I1' : { 'label' : 'Sơn nội thất Dulux Ambiance 5 trong 1 - bóng mờ', 'variant' : 'Nội thất', 'segment' : 'Siêu cao cấp', 'cities' : ['Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 9']},
    '_I1a' : { 'label' : 'Sơn nội thất Dulux Ambiance 5 trong 1 - siêu bóng', 'variant' : 'Nội thất', 'segment' : 'Siêu cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng'], 'groups' : ['Group 9']},
    '_I2' : { 'label' : 'Sơn nội thất Dulux Easy Clean kháng virus - mờ', 'variant' : 'Nội thất', 'segment' : 'Cao cấp', 'cities' : ['Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 9']},
    '_I2a' : { 'label' : 'Sơn nội thất Dulux Easy Clean kháng virus - bóng', 'variant' : 'Nội thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng'], 'groups' : ['Group 9']},
    '_I3' : { 'label' : 'Sơn nội thất Dulux Easy Clean (Lau chùi hiệu quả) - mờ', 'variant' : 'Nội thất', 'segment' : 'Cao cấp', 'cities' : ['Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 10']},
    '_I3a' : { 'label' : 'Sơn nội thất Dulux Easy Clean (Lau chùi hiệu quả) - bóng', 'variant' : 'Nội thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng'], 'groups' : ['Group 10']},
    '_I4a' : { 'label' : 'Sơn nội thất Dulux Inspire - bóng', 'variant' : 'Nội thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 12']},
    '_I4' : { 'label' : 'Sơn nội thất Dulux Inspire - mờ', 'variant' : 'Nội thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 12']},
    '_I5' : { 'label' : 'Sơn nội thất Maxilite Total - mờ', 'variant' : 'Nội thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 6', 'Group 13']},
    '_I6' : { 'label' : 'Sơn nội thất Maxilite Hi-Cover', 'variant' : 'Nội thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 6', 'Group 7', 'Group 13']},
    '_I7' : { 'label' : 'Sơn nội thất Maxilite Smooth', 'variant' : 'Nội thất', 'segment' : 'Gía rẻ', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 7']},
    '_I8' : { 'label' : 'Sơn nội thất JOTUN Majestic Dep Hoan Hao - mờ', 'variant' : 'Nội thất', 'segment' : 'Siêu cao cấp', 'cities' : ['Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 9']},
    '_I8a' : { 'label' : 'Sơn nội thất JOTUN Majestic Dep Hoan Hao - bóng', 'variant' : 'Nội thất', 'segment' : 'Siêu cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng'], 'groups' : ['Group 9']},
    '_I9' : { 'label' : 'Sơn nội thất Jotun Essence Dễ Lau Chùi ', 'variant' : 'Nội thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 10']},
    '_I10' : { 'label' : 'Sơn nội thất Jotun Jotaplast', 'variant' : 'Nội thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 6']},
    '_I11' : { 'label' : 'Sơn nội thất TOA 4 Seasons Top Silk', 'variant' : 'Nội thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 12']},
    '_I12' : { 'label' : 'Sơn nội thất Nippon Matex', 'variant' : 'Nội thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 6']},
    '_I13' : { 'label' : 'Sơn nội thất Nippon Matex Sắc Màu Dịu Mát', 'variant' : 'Nội thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 6']},
    '_I14' : { 'label' : 'Sơn nội thất Nippon Vatex', 'variant' : 'Nội thất', 'segment' : 'Gía rẻ', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 7']},
    '_I15' : { 'label' : 'Sơn nội thất Expo Easy', 'variant' : 'Nội thất', 'segment' : 'Gía rẻ', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 6']},
    '_E1' : { 'label' : 'Sơn ngoại thất cao cấp Dulux Weathershield', 'variant' : 'Ngoại thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 1']},
    '_E2' : { 'label' : 'Sơn ngoại thất siêu cao cấp Dulux Weathershield Powerflexx', 'variant' : 'Ngoại thất', 'segment' : 'Siêu cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 2']},
    '_E3' : { 'label' : 'Sơn ngoại thất cao cấp Dulux Weathershield Colour Protect', 'variant' : 'Ngoại thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 1']},
    '_E4' : { 'label' : 'Sơn ngoại thất Dulux Inspire Bóng/Mờ', 'variant' : 'Ngoại thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 5']},
    '_E4B' : { 'label' : 'Sơn ngoại thất Dulux Inspire Bóng', 'variant' : 'Ngoại thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 5']},
    '_E4M' : { 'label' : 'Sơn ngoại thất Dulux Inspire Mờ', 'variant' : 'Ngoại thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 5']},
    '_E5' : { 'label' : 'Sơn ngoại thất Maxilite Tough - mờ', 'variant' : 'Ngoại thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 8', 'Group 13']},
    '_E6' : { 'label' : 'Sơn ngoại thất Jotashield Chong Phai Mau', 'variant' : 'Ngoại thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 1']},
    '_E7' : { 'label' : 'Sơn ngoại thất Jotashield Ben Mau Toi Uu', 'variant' : 'Ngoại thất', 'segment' : 'Siêu cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 2']},
    '_E8' : { 'label' : 'Sơn ngoại thất Jotun Essence ngoại thất bền đẹp', 'variant' : 'Ngoại thất', 'segment' : 'Cận cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 5']},
    '_E9' : { 'label' : 'Sơn ngoại thất Jotun Jotatough', 'variant' : 'Ngoại thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 8']},
    '_E10' : { 'label' : 'Sơn ngoại thất Nippon Super Matex', 'variant' : 'Ngoại thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 8']},
    '_E11' : { 'label' : 'Sơn ngoại thất Expo Rainkote', 'variant' : 'Ngoại thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 8']},
    '_W1' : { 'label' : 'Chất Chống thấm tường Dulux Aquatech Flex', 'variant' : 'Chống thấm', 'segment' : 'Siêu cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 3']},
    '_W3' : { 'label' : 'Chất Chống thấm tường Dulux Aquatech Chống thấm vượt trội', 'variant' : 'Chống thấm', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Đà Nẵng'], 'groups' : ['Group 4']},
    '_W2' : { 'label' : 'Sơn Chống thấm tường Jotun Waterguard', 'variant' : 'Chống thấm', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 3']},
    '_W4' : { 'label' : 'Chất chống thấm tường Kova CT11A Gold', 'variant' : 'Chống thấm', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Đà Nẵng'], 'groups' : ['Group 4']},
    '_T1' : { 'label' : 'Sơn dầu Maxilite', 'variant' : 'Sơn dầu', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 11']},
    '_T2' : { 'label' : 'Sơn dầu Alkyld/ Sơn dầu dành cho gỗ & kim loại - Bạch Tuyết', 'variant' : 'Sơn dầu', 'segment' : 'Trung cấp', 'cities' : ['Hồ Chí Minh', 'Cần Thơ'], 'groups' : ['Group 11']},
    '_T3' : { 'label' : 'Sơn dầu dành cho gỗ & kim loại - Lobster', 'variant' : 'Sơn dầu', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 11']},
    '_I5e' : { 'label' : 'Sơn nội thất Maxilite Total - mờ (bao bì mới)', 'variant' : 'Nội thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 13']},
    '_I6e' : { 'label' : 'Sơn nội thất Maxilite Hi-Cover - mờ (bao bì mới)', 'variant' : 'Nội thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 13']},
    '_E5e' : { 'label' : 'Sơn ngoại thất Maxilite Tough - mờ (bao bì mới)', 'variant' : 'Ngoại thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 13']},
    '_E12e' : { 'label' : 'Sơn ngoại thất Jotun Tough Shield - mờ', 'variant' : 'Ngoại thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 13']},
    '_E13e' : { 'label' : 'Sơn ngoại thất Jotun Tough Shield Max - mờ', 'variant' : 'Ngoại thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội', 'Hải Phòng', 'Hồ Chí Minh', 'Cần Thơ', 'Đà Nẵng'], 'groups' : ['Group 13']}
}

cities = ["Hồ Chí Minh","Hà Nội", "Hải Phòng","Đà Nẵng", "Cần Thơ"]
segments = ["Siêu cao cấp", "Cao cấp", "Cận cao cấp", "Trung cấp", "Giá rẻ"]

writer = pd.ExcelWriter("{}/output_2.xlsx".format(project_name), engine='xlsxwriter')

cols = ["_Q1"]
cols.extend(list(filter(lambda x: re.match("(_Phase\[{_1}\])(.*)((_E4)|(_Q5_((SKU)|(GiaTien))))$", x), df.columns)))

df_data = df.loc[:, cols]
df_data.reset_index(inplace=True)
df_data.set_index(["Respondent.ID", "_Q1"], inplace=True)

column_names = ["Q5_SKU", "Q5_PRICE"]

for c in column_names:
    df_data[c] = df_data[df_data.columns[0:6]].apply(lambda x: ','.join(x.astype(str)), axis=1)
    df_data.drop(df_data.columns[0:6], axis=1, inplace=True)

df_data.loc[df_data["_Phase[{_1}]._E4"] == 11, "Q5_SKU"] = df_data["Q5_SKU"].str.replace("_E4", "_E4B")
df_data.loc[df_data["_Phase[{_1}]._E4"] == 12, "Q5_SKU"] = df_data["Q5_SKU"].str.replace("_E4", "_E4M")

df_q5 = df_data.stack().str.split(',', expand=True).stack().unstack(-2).reset_index()
df_q5.drop("level_2", inplace=True, axis=1)
df_q5 = df_q5[df_q5["Q5_PRICE"] != 'nan']
df_q5["Q5_PRICE"] = df_q5["Q5_PRICE"].apply(lambda x: float(x))
df_q5.set_index(["Respondent.ID", "_Q1", "Q5_SKU"], inplace=True)

df_q5.to_excel(writer, sheet_name="Q5")
worksheet = writer.sheets["Q5"]

cols = ["_Q1"]
cols.extend(list(filter(lambda x: re.match("(_Phase\[{_1}\])(.*)((_E4)|(_Q5_7)|(_Q5b_SKU)|(_Q5b_GiaTien))$", x), df.columns)))

df_data = df.loc[:, cols]
df_data.reset_index(inplace=True)
df_data.set_index(["Respondent.ID", "_Q1"], inplace=True)

df_data = df_data[df_data["_Phase[{_1}]._Q5_7"] == 1]
df_data.drop("_Phase[{_1}]._Q5_7", inplace=True, axis=1)

column_names = ["Q5_SKU", "Q5b_PRICE"]

for c in column_names:
    df_data[c] = df_data[df_data.columns[0:11]].apply(lambda x: ','.join(x.astype(str)), axis=1)
    df_data.drop(df_data.columns[0:11], axis=1, inplace=True)

df_data.loc[df_data["_Phase[{_1}]._E4"] == 11, "Q5_SKU"] = df_data["Q5_SKU"].str.replace("_E4", "_E4B")
df_data.loc[df_data["_Phase[{_1}]._E4"] == 12, "Q5_SKU"] = df_data["Q5_SKU"].str.replace("_E4", "_E4M")

df_q5b = df_data.stack().str.split(',', expand=True).stack().unstack(-2).reset_index()
df_q5b.drop("level_2", inplace=True, axis=1)
df_q5b = df_q5b[df_q5b["Q5b_PRICE"] != 'nan']
df_q5b["Q5b_PRICE"] = df_q5b["Q5b_PRICE"].apply(lambda x: float(x))
df_q5b = df_q5b[df_q5b["Q5b_PRICE"] > 0]
df_q5b.set_index(["Respondent.ID", "_Q1", "Q5_SKU"], inplace=True)

df_merge = pd.merge(df_q5, df_q5b, left_index=True, right_index=True, how="outer")
df_merge.reset_index(inplace=True)

df_q5b.to_excel(writer, sheet_name="Q5B")
worksheet = writer.sheets["Q5B"]

df_merge["PRICE"] = df_merge.apply(lambda x: x['Q5_PRICE'] if x['Q5_PRICE'] > 0 else x['Q5b_PRICE'], axis=1)

for k in obj_skus_define.keys():
    df_merge.loc[df_merge["Q5_SKU"] == k, "GROUPS"] = ','.join(obj_skus_define[k]['groups'])

df_merge.set_index(["Respondent.ID","_Q1","Q5_SKU","Q5_PRICE","Q5b_PRICE","PRICE"], inplace=True)

df_merge = df_merge.stack().str.split(',', expand=True).stack().reset_index().rename(columns={0 : "GROUP"})
df_merge.drop(["level_6","level_7"], axis=1, inplace=True)

df_merge.to_excel(writer, sheet_name="Q5B_MERGE")
worksheet = writer.sheets["Q5B_MERGE"]

df_count = df_merge[["Respondent.ID","_Q1","Q5_SKU"]].groupby(["_Q1","Q5_SKU"], as_index=True).count().reset_index(-2, drop=False).reset_index(-1, drop=False)

obj_counts = {}

for i, s in obj_skus_define.items():
    for j in cities:
        if j not in obj_counts.keys():
            obj_counts[j] = {}
        if j in obj_skus_define[i]['cities']:
            if i not in obj_counts[j].keys():
                obj_counts[j][i] = {}
            try:
                obj_counts[j][i] = df_count.loc[(df_count["_Q1"] == j) & (df_count["Q5_SKU"] == i)]["Respondent.ID"].values[0]
            except IndexError:
                obj_counts[j][i] = 0

for i, c in obj_counts.items():
    obj_city = {'sku_name' : [], 'ideal_sample' : [], 'current_sample' : [], 'remaining_sample' : []}
    
    if (i == "Hồ Chí Minh" or i == "Hà Nội"):
        base = 15
    elif i == "Hải Phòng":
        base = 7
    elif i == "Đà Nẵng":
        base = 11
    else:
        base = 8

    for j, b in c.items():
        obj_city['sku_name'].append(j)
        obj_city['ideal_sample'].append(b)
        obj_city['current_sample'].append(base)
        obj_city['remaining_sample'].append(base - b)
    
    df_4 = pd.DataFrame(data=obj_city)
    df_4.to_excel(writer, sheet_name=i)
    worksheet = writer.sheets[i]

writer.save()

presentation = Presentation()
title_slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "TABLO - SKU PRICE BY CITY"
subtitle.text = "Python was here!"

fig, ax = plt.subplots(nrows=1, ncols=1, figsize = (10, 6))

for s in obj_skus_define.keys():
    df_merge_filter = df_merge[df_merge["Q5_SKU"] == s]
    
    if not df_merge_filter.empty:
        
        ax.ticklabel_format(style='plain')
            
        boxplot = df_merge_filter.boxplot(column=["PRICE"], by="_Q1", ax=ax, rot=0, grid=True, showfliers=True, showmeans=True, meanline=True, sym="g o")
        
        boxplot.set_title("SKU {} - {}".format(s[1:len(s)], obj_skus_define[s]['label']))
        boxplot.set_xlabel('City')
        boxplot.set_ylabel('Price')

        pic_name = 'plot{}.png'.format(s)
        plt.savefig(r'{}/images/{}'.format(project_name, pic_name))
        plt.cla()
        
        blank_slide_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(blank_slide_layout)
        shapes = slide.shapes
        
        left = Inches(0)
        top = Inches(1)
        
        pic = shapes.add_picture(r'{}/images/{}'.format(project_name, pic_name), left, top)

pptx_name = 'python - tablo.pptx'
presentation.save('{}/{}'.format(project_name, pptx_name))

"""

presentation = Presentation()
title_slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "TABLO - SKU PRICE BY CITY"
subtitle.text = "Python was here!"

obj_skus = {}

for i, r in df_data.iterrows():
    for i in range(1, 7):
        if math.isnan(r[i + 6]) is False:
            if r[i] not in obj_skus.keys():
                obj_skus[r[i]] = {}
            if r[0] not in obj_skus[r[i]].keys():
                obj_skus[r[i]][r[0]] = list()
            
            obj_skus[r[i]][r[0]].append(r[i + 6])

fig, ax = plt.subplots(nrows=1, ncols=1, figsize = (10, 6))

for i, s in obj_skus.items():    
    print(s.values())
    ax.boxplot(s.values(), showfliers=True, showmeans=True, meanline=True, sym="g o")
    #ax.boxplot(s.values(), notch=False, showfliers=True, sym='+', vert=True, whis=1.5, positions=None, widths=None, patch_artist=False, bootstrap=None, usermedians=None, conf_intervals=None)
    ax.set_xticklabels(s.keys())
    ax.set_title("{}".format(i))

    pic_name = 'plot{}.png'.format(i)

    plt.savefig(r'{}/images/{}'.format(project_name,pic_name))
    plt.cla()

    blank_slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(blank_slide_layout)
    shapes = slide.shapes

    shapes.title.text = "SKU {}".format(i[1:len(i)])

    left = Inches(0)
    top = Inches(1.5)
    
    pic = shapes.add_picture(r'{}/images/{}'.format(project_name, pic_name), left, top)

title_slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "TABLO - SKU PRICE BY SKU"
subtitle.text = "Python was here!"

obj_cities = {}

for i, r in df_data.iterrows():
    if r["_Q1"] not in obj_cities.keys():
        obj_cities[r["_Q1"]] = {}
    for ic in range(1, 7):
        if math.isnan(r[ic + 6]) is False:
            segment = obj_skus_define[r[ic]]['segment']

            if segment not in obj_cities[r["_Q1"]]:
                obj_cities[r["_Q1"]][segment] = {}

            if r[ic] not in obj_cities[r["_Q1"]][segment]:
                obj_cities[r["_Q1"]][segment][r[ic]] = list()
            
            obj_cities[r["_Q1"]][segment][r[ic]].append(r[ic + 6])

for i, c in obj_cities.items():
    for j, s in c.items():
        print(s.values())
        ax.boxplot(s.values(), showfliers=True, showmeans=True, meanline=True, sym="g o") 
        #ax.boxplot(s.values(), notch=False, showfliers=True, sym='+', vert=True, whis=1.5, positions=None, widths=None, patch_artist=False, bootstrap=None, usermedians=None, conf_intervals=None)
        ax.set_xticklabels(s.keys())
        ax.set_title("{} - {}".format(i, j))

        pic_name = 'plot{}.png'.format(i)

        plt.savefig(r'{}/images/{}'.format(project_name, pic_name))
        plt.cla()

        blank_slide_layout = presentation.slide_layouts[5]
        slide = presentation.slides.add_slide(blank_slide_layout)
        shapes = slide.shapes

        shapes.title.text = "SKU {} {}".format(i, j)

        left = Inches(0)
        top = Inches(1.5)
        
        pic = shapes.add_picture(r'{}/images/{}'.format(project_name, pic_name), left, top)

pptx_name = 'python - tablo.pptx'

presentation.save('{}/{}'.format(project_name, pptx_name))

"""




"""
fig, axs = plt.subplots(nrows=len(cities), ncols=len(segments), figsize = (15, 10))
fig.subplots_adjust(left=0.08, right=0.98, bottom=0.05, top=0.9, hspace=0.4, wspace=0.3)

_col, _row = 0, 0

for c in cities:
    _col = 0

    for s in segments:
        df_cities = df_data.query("_Q1 == '{}'".format(c))

        obj_skus = {}

        for index, row in df_cities.iterrows():
            for i in range(1, 7):
                if math.isnan(row[i + 6]) is False:
                    if obj_skus_define[row[i]]['segment'] == s:
                        if row[i] not in obj_skus.keys():
                            obj_skus[row[i]] = list()
                        obj_skus[row[i]].append(row[i + 6]) 

        d1 = []

        for i in obj_skus:
            if len(obj_skus[i]) > 0:
                d1.append(obj_skus[i])
       
        axs[_row, _col].boxplot(d1, showmeans=True, meanline=True, sym="g o")
        axs[_row, _col].set_xticklabels(obj_skus.keys())
        axs[_row, _col].set_title("{} - {}".format(c, s))
        _col += 1 
    _row += 1

plt.savefig('plot.png')
writer = pd.ExcelWriter("output.xlsx", engine='xlsxwriter')

df.to_excel(writer, sheet_name="Output")
worksheet = writer.sheets["Output"]
worksheet.insert_image('A1', 'plot.png')
writer.save()

plt.show()

"""
