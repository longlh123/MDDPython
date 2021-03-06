from email.policy import default
from inspect import stack
from ipaddress import v4_int_to_packed
import math
from re import T
from turtle import left, pos
from xml.dom.pulldom import ErrorHandler
import win32com.client as w32
from metadata import Metadata
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import openpyxl
from pptx import Presentation
from pptx.util import Inches

#MDM = w32.Dispatch('MDM.Document')
#MDM.Open(r'metadatas/PVN2021249_F61_W1_2022v1.mdd')

#questions = ["InstanceID","_Year","_Month","_Quarter","_ResProvinces","_City_Groups","_Target","_Users","_Class","_EPI_OOP_Inject","_Table_1","_Table_2","_Table_8","_Table_9","_Table_10_Before","_Table_10","_Table_13","_Table_14","_Table_15","_Table_16","_Table_17","_Table_18_Before","_Table_18","_Table_19","_Table_20","_Table_21","_Table_3","_UA1","_UA2","_UA3a","_UA3b","_UB1","_UB2","_UB4","_UB5","_Q3","_S16","_S16b","_S16c","_Q17","_UD3","_OP5c","_OP5e","_OP5d"]

questions = ["Respondent.ID","_Q1","_Phase"]

#try:
m = Metadata(r'metadatas/S22015945_DATA.mdd', r'metadatas/S22015945_DATA.ddf', questions)

df = m.convertToDataFrame() 
df = df.set_index("Respondent.ID")

obj_skus_define = {
    '_I1' : { 'variant' : 'Nội thất', 'segment' : 'Siêu cao cấp', 'cities' : ['Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_I1a' : { 'variant' : 'Nội thất', 'segment' : 'Siêu cao cấp', 'cities' : ['Hà Nội','Hải Phòng']},
    '_I2' : { 'variant' : 'Nội thất', 'segment' : 'Cao cấp', 'cities' : ['Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_I2a' : { 'variant' : 'Nội thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội','Hải Phòng']},
    '_I3' : { 'variant' : 'Nội thất', 'segment' : 'Cao cấp', 'cities' : ['Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_I3a' : { 'variant' : 'Nội thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội','Hải Phòng']},
    '_I4a' : { 'variant' : 'Nội thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_I4' : { 'variant' : 'Nội thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_I5' : { 'variant' : 'Nội thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_I6' : { 'variant' : 'Nội thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_I7' : { 'variant' : 'Nội thất', 'segment' : 'Gía rẻ', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_I8' : { 'variant' : 'Nội thất', 'segment' : 'Siêu cao cấp', 'cities' : ['Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_I8a' : { 'variant' : 'Nội thất', 'segment' : 'Siêu cao cấp', 'cities' : ['Hà Nội','Hải Phòng']},
    '_I9' : { 'variant' : 'Nội thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_I10' : { 'variant' : 'Nội thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_I11' : { 'variant' : 'Nội thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_I12' : { 'variant' : 'Nội thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_I13' : { 'variant' : 'Nội thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_I14' : { 'variant' : 'Nội thất', 'segment' : 'Gía rẻ', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_I15' : { 'variant' : 'Nội thất', 'segment' : 'Gía rẻ', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_E1' : { 'variant' : 'Ngoại thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_E2' : { 'variant' : 'Ngoại thất', 'segment' : 'Siêu cao cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_E3' : { 'variant' : 'Ngoại thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_E4' : { 'variant' : 'Ngoại thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']}, #]4B
    '_E4M' : { 'variant' : 'Ngoại thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_E5' : { 'variant' : 'Ngoại thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_E6' : { 'variant' : 'Ngoại thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_E7' : { 'variant' : 'Ngoại thất', 'segment' : 'Siêu cao cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_E8' : { 'variant' : 'Ngoại thất', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_E9' : { 'variant' : 'Ngoại thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_E10' : { 'variant' : 'Ngoại thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_E11' : { 'variant' : 'Ngoại thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_W1' : { 'variant' : 'Chống thấm', 'segment' : 'Siêu cao cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_W3' : { 'variant' : 'Chống thấm', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Đà Nẵng']},
    '_W2' : { 'variant' : 'Chống thấm', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_W4' : { 'variant' : 'Chống thấm', 'segment' : 'Cao cấp', 'cities' : ['Hà Nội','Hải Phòng','Đà Nẵng']},
    '_T1' : { 'variant' : 'Sơn dầu', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội','Hải Phòng','Cần Thơ','Đà Nẵng']},
    '_T2' : { 'variant' : 'Sơn dầu', 'segment' : 'Trung cấp', 'cities' : ['Hồ Chí Minh','Cần Thơ']},
    '_T3' : { 'variant' : 'Sơn dầu', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_I5e' : { 'variant' : 'Nội thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_I6e' : { 'variant' : 'Nội thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_E5e' : { 'variant' : 'Ngoại thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_E12e' : { 'variant' : 'Ngoại thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']},
    '_E13e' : { 'variant' : 'Ngoại thất', 'segment' : 'Trung cấp', 'cities' : ['Hà Nội','Hải Phòng','Hồ Chí Minh','Cần Thơ','Đà Nẵng']}
}

cities = ["Hồ Chí Minh","Hà Nội", "Hải Phòng","Đà Nẵng", "Cần Thơ"]
segments = ["Siêu cao cấp", "Cao cấp", "Cận cao cấp", "Trung cấp", "Giá rẻ"]

df_data = df.loc[:, ["_Q1","_Phase[{_1}]._Q5[{_1}]._Q5_SKU","_Phase[{_1}]._Q5[{_2}]._Q5_SKU","_Phase[{_1}]._Q5[{_3}]._Q5_SKU","_Phase[{_1}]._Q5[{_4}]._Q5_SKU","_Phase[{_1}]._Q5[{_5}]._Q5_SKU","_Phase[{_1}]._Q5[{_6}]._Q5_SKU","_Phase[{_1}]._Q5[{_1}]._Q5_GiaTien","_Phase[{_1}]._Q5[{_2}]._Q5_GiaTien","_Phase[{_1}]._Q5[{_3}]._Q5_GiaTien","_Phase[{_1}]._Q5[{_4}]._Q5_GiaTien","_Phase[{_1}]._Q5[{_5}]._Q5_GiaTien","_Phase[{_1}]._Q5[{_6}]._Q5_GiaTien"]]

#BAR CHART - 100% STACKED COLUMNS CHART

#Combine multiple columns into a single dataframe 
df_data["Q5_SKU"] = df_data[df_data.columns[1:7]].apply(lambda x: ','.join(x.astype(str)), axis=1)
df_data["Q5_PRICE"] = df_data[df_data.columns[7:13]].apply(lambda x: ','.join(x.astype(str)), axis=1)

df_data.drop(df_data.columns[1:13], axis=1, inplace=True)

df_data.reset_index(-1, drop=False, inplace=True)
df_data.set_index(["Respondent.ID","_Q1"], drop=True, inplace=True, verify_integrity=False)

df2 = df_data.stack().str.split(',', expand=True).stack().unstack(-2).reset_index(-1, drop=True).reset_index()

df_3 = df2[["Respondent.ID","_Q1","Q5_SKU"]].groupby(["_Q1","Q5_SKU"], as_index=True).count().reset_index(-2, drop=False).reset_index(-1, drop=False)

obj_counts = {}

for i, s in obj_skus_define.items():
    for j in cities:
        if j not in obj_counts.keys():
            obj_counts[j] = {}
        if j in obj_skus_define[i]['cities']:
            if i not in obj_counts[j].keys():
                obj_counts[j][i] = {}
            try:
                obj_counts[j][i] = df_3.loc[(df_3["_Q1"] == j) & (df_3["Q5_SKU"] == i)]["Respondent.ID"].values[0]
            except IndexError:
                obj_counts[j][i] = 0

writer = pd.ExcelWriter("output_2.xlsx", engine='xlsxwriter')

for i, c in obj_counts.items():
    obj_city = {'sku_name' : [], 'ideal_sample' : [], 'current_sample' : [], 'remaining_sample' : []}
    
    if (i == "Hồ Chí Minh" or i == "Hà Nội"):
        base = 15
    elif i == "Hải Phòng":
        base = 7
    elif i == "Đà Nẵng":
        base = 11
    else:
        base = 8

    for j, b in c.items():
        obj_city['sku_name'].append(j)
        obj_city['ideal_sample'].append(b)
        obj_city['current_sample'].append(base)
        obj_city['remaining_sample'].append(base - b)
    
    df_4 = pd.DataFrame(data=obj_city)
    df_4.to_excel(writer, sheet_name=i)
    worksheet = writer.sheets[i]

writer.save()




#

#print(df_4)






#f_3.loc[df_3["_Q1"] == "Hồ Chí Minh", "SKU_BASE"] = 15 
##df_3.loc[df_3["_Q1"] == "Cần Thơ", "SKU_BASE"] = 8 
##df_3.loc[df_3["_Q1"] == "Hà Nội", "SKU_BASE"] = 15 
#df_3.loc[df_3["_Q1"] == "Hải Phòng", "SKU_BASE"] = 7 
#df_3.loc[df_3["_Q1"] == "Đà Nẵng", "SKU_BASE"] = 11 

#df_3["SKU_TEMPT"] = df_3["SKU_BASE"] - df_3["Respondent.ID"]

#writer = pd.ExcelWriter("output_2.xlsx", engine='xlsxwriter')

##df_3.to_excel(writer, sheet_name="Output")
#worksheet = writer.sheets["Output"]
#writer.save()



"""

presentation = Presentation()
title_slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "TABLO - SKU PRICE BY CITY"
subtitle.text = "Python was here!"

obj_skus = {}

for i, r in df_data.iterrows():
    for i in range(1, 7):
        if math.isnan(r[i + 6]) is False:
            if r[i] not in obj_skus.keys():
                obj_skus[r[i]] = {}
            if r[0] not in obj_skus[r[i]].keys():
                obj_skus[r[i]][r[0]] = list()
            
            obj_skus[r[i]][r[0]].append(r[i + 6])

fig, ax = plt.subplots(nrows=1, ncols=1, figsize = (10, 6))

for i, s in obj_skus.items():    
    print(s.values())
    ax.boxplot(s.values(), showfliers=True, showmeans=True, meanline=True, sym="g o")
    #ax.boxplot(s.values(), notch=False, showfliers=True, sym='+', vert=True, whis=1.5, positions=None, widths=None, patch_artist=False, bootstrap=None, usermedians=None, conf_intervals=None)
    ax.set_xticklabels(s.keys())
    ax.set_title("{}".format(i))

    plt.savefig('plot{}.png'.format(i))
    plt.cla()

    blank_slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(blank_slide_layout)
    shapes = slide.shapes

    shapes.title.text = "SKU {}".format(i[1:len(i)])

    left = Inches(0)
    top = Inches(1.5)
    
    pic = shapes.add_picture('plot{}.png'.format(i), left, top)

title_slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "TABLO - SKU PRICE BY SKU"
subtitle.text = "Python was here!"

obj_cities = {}

for i, r in df_data.iterrows():
    if r["_Q1"] not in obj_cities.keys():
        obj_cities[r["_Q1"]] = {}
    for ic in range(1, 7):
        if math.isnan(r[ic + 6]) is False:
            segment = obj_skus_define[r[ic]]['segment']

            if segment not in obj_cities[r["_Q1"]]:
                obj_cities[r["_Q1"]][segment] = {}

            if r[ic] not in obj_cities[r["_Q1"]][segment]:
                obj_cities[r["_Q1"]][segment][r[ic]] = list()
            
            obj_cities[r["_Q1"]][segment][r[ic]].append(r[ic + 6])

for i, c in obj_cities.items():
    for j, s in c.items():
        print(s.values())
        ax.boxplot(s.values(), showfliers=True, showmeans=True, meanline=True, sym="g o") 
        #ax.boxplot(s.values(), notch=False, showfliers=True, sym='+', vert=True, whis=1.5, positions=None, widths=None, patch_artist=False, bootstrap=None, usermedians=None, conf_intervals=None)
        ax.set_xticklabels(s.keys())
        ax.set_title("{} - {}".format(i, j))

        plt.savefig('plot{}{}.png'.format(i, j))
        plt.cla()

        blank_slide_layout = presentation.slide_layouts[5]
        slide = presentation.slides.add_slide(blank_slide_layout)
        shapes = slide.shapes

        shapes.title.text = "SKU {} {}".format(i, j)

        left = Inches(0)
        top = Inches(1.5)
        
        pic = shapes.add_picture('plot{}{}.png'.format(i, j), left, top)


presentation.save('python - tablo.pptx')

"""


"""
fig, axs = plt.subplots(nrows=len(cities), ncols=len(segments), figsize = (15, 10))
fig.subplots_adjust(left=0.08, right=0.98, bottom=0.05, top=0.9, hspace=0.4, wspace=0.3)

_col, _row = 0, 0

for c in cities:
    _col = 0

    for s in segments:
        df_cities = df_data.query("_Q1 == '{}'".format(c))

        obj_skus = {}

        for index, row in df_cities.iterrows():
            for i in range(1, 7):
                if math.isnan(row[i + 6]) is False:
                    if obj_skus_define[row[i]]['segment'] == s:
                        if row[i] not in obj_skus.keys():
                            obj_skus[row[i]] = list()
                        obj_skus[row[i]].append(row[i + 6]) 

        d1 = []

        for i in obj_skus:
            if len(obj_skus[i]) > 0:
                d1.append(obj_skus[i])
       
        axs[_row, _col].boxplot(d1, showmeans=True, meanline=True, sym="g o")
        axs[_row, _col].set_xticklabels(obj_skus.keys())
        axs[_row, _col].set_title("{} - {}".format(c, s))
        _col += 1 
    _row += 1

plt.savefig('plot.png')
writer = pd.ExcelWriter("output.xlsx", engine='xlsxwriter')

df.to_excel(writer, sheet_name="Output")
worksheet = writer.sheets["Output"]
worksheet.insert_image('A1', 'plot.png')
writer.save()

plt.show()

"""
